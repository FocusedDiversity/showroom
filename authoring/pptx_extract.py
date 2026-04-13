"""
PPTX Extractor — extracts colors, fonts, style cues, and images from a .pptx file.

Uses python-pptx to read the slide master theme and produce data compatible
with the Palette and Layout dataclasses.
"""

import re
from lxml import etree
from pptx import Presentation
from pptx.util import Emu


# ── Color Extraction ─────────────────────────────────────────────────

def _hex_from_rgb(r, g, b):
    """Convert RGB ints to hex string."""
    return f'#{r:02x}{g:02x}{b:02x}'


def _lighten(hex_color, factor=0.4):
    """Lighten a hex color by blending toward white."""
    hex_color = hex_color.lstrip('#')
    r, g, b = int(hex_color[:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
    r = int(r + (255 - r) * factor)
    g = int(g + (255 - g) * factor)
    b = int(b + (255 - b) * factor)
    return _hex_from_rgb(r, g, b)


def _darken(hex_color, factor=0.3):
    """Darken a hex color by blending toward black."""
    hex_color = hex_color.lstrip('#')
    r, g, b = int(hex_color[:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
    r = int(r * (1 - factor))
    g = int(g * (1 - factor))
    b = int(b * (1 - factor))
    return _hex_from_rgb(r, g, b)


def _hex_to_rgba(hex_color, alpha):
    """Convert hex color to rgba() CSS string."""
    hex_color = hex_color.lstrip('#')
    r, g, b = int(hex_color[:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
    return f'rgba({r},{g},{b},{alpha})'


def _extract_theme_colors(prs):
    """Extract the theme color scheme from the first slide master.

    Returns a dict of theme color names → hex values.
    """
    ns = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    }

    colors = {}

    try:
        master = prs.slide_masters[0]
        # Navigate to the theme element
        theme_el = master.element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}theme')
        if theme_el is None:
            # Try via the slide master's relationships
            for rel in master.part.rels.values():
                if 'theme' in rel.reltype:
                    theme_part = rel.target_part
                    theme_el = etree.fromstring(theme_part.blob)
                    break

        if theme_el is None:
            return colors

        # Find clrScheme
        clr_scheme = theme_el.find('.//a:clrScheme', ns)
        if clr_scheme is None:
            return colors

        # Extract each color slot
        color_names = ['dk1', 'dk2', 'lt1', 'lt2',
                       'accent1', 'accent2', 'accent3', 'accent4',
                       'accent5', 'accent6', 'hlink', 'folHlink']

        for name in color_names:
            el = clr_scheme.find(f'a:{name}', ns)
            if el is None:
                continue

            # Color can be srgbClr or sysClr
            srgb = el.find('a:srgbClr', ns)
            sys_clr = el.find('a:sysClr', ns)

            if srgb is not None:
                colors[name] = f'#{srgb.get("val")}'
            elif sys_clr is not None:
                last_clr = sys_clr.get('lastClr')
                if last_clr:
                    colors[name] = f'#{last_clr}'
                else:
                    # Map system color names to reasonable defaults
                    sys_map = {
                        'windowText': '#000000',
                        'window': '#ffffff',
                    }
                    colors[name] = sys_map.get(sys_clr.get('val'), '#888888')

    except (IndexError, AttributeError):
        pass

    return colors


def extract_colors(pptx_path):
    """Extract colors from a PPTX and return a dict of Palette constructor kwargs.

    Args:
        pptx_path: Path to the .pptx file

    Returns:
        dict with all fields needed to construct a Palette
    """
    prs = Presentation(pptx_path)
    theme_colors = _extract_theme_colors(prs)

    # Defaults (Office default theme)
    defaults = {
        'dk1': '#000000', 'dk2': '#44546A', 'lt1': '#FFFFFF', 'lt2': '#E7E6E6',
        'accent1': '#4472C4', 'accent2': '#ED7D31', 'accent3': '#A5A5A5',
        'accent4': '#FFC000', 'accent5': '#5B9BD5', 'accent6': '#70AD47',
    }
    c = {**defaults, **theme_colors}

    # Map theme colors to Palette fields
    accent_primary = c['accent1']
    accent_secondary = c['accent2']
    accent_tertiary = c['accent3']
    bg_dark = c['dk1'] if c['dk1'] != '#000000' else _darken(c['dk2'], 0.4)
    bg_light = c['lt1'] if c['lt1'] != '#ffffff' else '#FAF8F5'
    text_dark = c['dk1'] if c['dk1'] != '#000000' else c['dk2']
    text_light = c['lt1']
    text_muted = _lighten(text_dark, 0.4)

    return {
        'background_dark': bg_dark,
        'background_light': bg_light,
        'accent_primary': accent_primary,
        'accent_secondary': accent_secondary,
        'accent_tertiary': accent_tertiary,
        'text_dark': text_dark,
        'text_light': text_light,
        'text_muted': text_muted,
        'top_bar_gradient': f'linear-gradient(90deg, {accent_primary}, {accent_secondary}, {accent_tertiary})',
        'card_border_light': _lighten(accent_primary, 0.6),
        'card_border_dark': _hex_to_rgba(accent_primary, 0.12),
        'card_bg_light': '#ffffff',
        'card_bg_dark': _hex_to_rgba(accent_primary, 0.10),
        'section_label_light': accent_secondary,
        'section_label_dark': accent_tertiary,
        'highlight_light': accent_primary,
        'highlight_dark_gradient': f'linear-gradient(135deg, {accent_tertiary}, {accent_primary})',
    }


# ── Font Extraction ──────────────────────────────────────────────────

# Common PPTX font → Google Fonts mapping
GOOGLE_FONTS_MAP = {
    # Sans-serif
    'calibri': {'family': "'Open Sans', sans-serif", 'import': 'Open+Sans:wght@300;400;500;600;700', 'category': 'sans-serif'},
    'arial': {'family': "'Open Sans', sans-serif", 'import': 'Open+Sans:wght@300;400;500;600;700', 'category': 'sans-serif'},
    'helvetica': {'family': "'Inter', sans-serif", 'import': 'Inter:wght@300;400;500;600;700', 'category': 'sans-serif'},
    'helvetica neue': {'family': "'Inter', sans-serif", 'import': 'Inter:wght@300;400;500;600;700', 'category': 'sans-serif'},
    'century gothic': {'family': "'Quicksand', sans-serif", 'import': 'Quicksand:wght@300;400;500;600;700', 'category': 'sans-serif'},
    'trebuchet ms': {'family': "'Source Sans 3', sans-serif", 'import': 'Source+Sans+3:wght@300;400;500;600;700', 'category': 'sans-serif'},
    'verdana': {'family': "'Nunito', sans-serif", 'import': 'Nunito:wght@300;400;500;600;700', 'category': 'sans-serif'},
    'tahoma': {'family': "'Open Sans', sans-serif", 'import': 'Open+Sans:wght@300;400;500;600;700', 'category': 'sans-serif'},
    'segoe ui': {'family': "'Inter', sans-serif", 'import': 'Inter:wght@300;400;500;600;700', 'category': 'sans-serif'},
    'gill sans': {'family': "'Lato', sans-serif", 'import': 'Lato:wght@300;400;700;900', 'category': 'sans-serif'},
    'gill sans mt': {'family': "'Lato', sans-serif", 'import': 'Lato:wght@300;400;700;900', 'category': 'sans-serif'},
    'avenir': {'family': "'Nunito Sans', sans-serif", 'import': 'Nunito+Sans:wght@300;400;500;600;700', 'category': 'sans-serif'},
    'futura': {'family': "'Poppins', sans-serif", 'import': 'Poppins:wght@300;400;500;600;700', 'category': 'sans-serif'},
    'proxima nova': {'family': "'Montserrat', sans-serif", 'import': 'Montserrat:wght@300;400;500;600;700', 'category': 'sans-serif'},
    'myriad pro': {'family': "'Source Sans 3', sans-serif", 'import': 'Source+Sans+3:wght@300;400;500;600;700', 'category': 'sans-serif'},
    'franklin gothic': {'family': "'Libre Franklin', sans-serif", 'import': 'Libre+Franklin:wght@300;400;500;600;700', 'category': 'sans-serif'},
    'gotham': {'family': "'Montserrat', sans-serif", 'import': 'Montserrat:wght@300;400;500;600;700', 'category': 'sans-serif'},
    'brandon grotesque': {'family': "'Nunito', sans-serif", 'import': 'Nunito:wght@300;400;500;600;700', 'category': 'sans-serif'},

    # Serif
    'garamond': {'family': "'EB Garamond', serif", 'import': 'EB+Garamond:wght@400;500;600;700', 'category': 'serif'},
    'georgia': {'family': "'Lora', serif", 'import': 'Lora:wght@400;500;600;700', 'category': 'serif'},
    'times new roman': {'family': "'Playfair Display', serif", 'import': 'Playfair+Display:wght@400;500;600;700', 'category': 'serif'},
    'cambria': {'family': "'Merriweather', serif", 'import': 'Merriweather:wght@300;400;700;900', 'category': 'serif'},
    'palatino': {'family': "'Libre Baskerville', serif", 'import': 'Libre+Baskerville:wght@400;700', 'category': 'serif'},
    'palatino linotype': {'family': "'Libre Baskerville', serif", 'import': 'Libre+Baskerville:wght@400;700', 'category': 'serif'},
    'book antiqua': {'family': "'Libre Baskerville', serif", 'import': 'Libre+Baskerville:wght@400;700', 'category': 'serif'},
    'rockwell': {'family': "'Zilla Slab', serif", 'import': 'Zilla+Slab:wght@300;400;500;600;700', 'category': 'serif'},
    'century': {'family': "'EB Garamond', serif", 'import': 'EB+Garamond:wght@400;500;600;700', 'category': 'serif'},
    'baskerville': {'family': "'Libre Baskerville', serif", 'import': 'Libre+Baskerville:wght@400;700', 'category': 'serif'},

    # Display
    'impact': {'family': "'Oswald', sans-serif", 'import': 'Oswald:wght@400;500;600;700', 'category': 'display'},
    'copperplate': {'family': "'Cinzel', serif", 'import': 'Cinzel:wght@400;500;600;700', 'category': 'display'},

    # Synaptiq brand fonts (already in the system)
    'museo slab': {'family': "'Zilla Slab', serif", 'import': 'Zilla+Slab:wght@300;400;500;600;700', 'category': 'serif'},
    'arboria': {'family': "'Quicksand', sans-serif", 'import': 'Quicksand:wght@300;400;500;600;700', 'category': 'sans-serif'},
    'all round gothic': {'family': "'Quicksand', sans-serif", 'import': 'Quicksand:wght@300;400;500;600;700', 'category': 'sans-serif'},
}

# Default fallback
_DEFAULT_FONT = {'family': "'Quicksand', sans-serif", 'import': 'Quicksand:wght@300;400;500;600;700', 'category': 'sans-serif'}


def map_to_google_fonts(font_name):
    """Map a PPTX font name to its closest Google Fonts equivalent.

    Returns dict with 'family', 'import', 'category'.
    """
    if not font_name:
        return _DEFAULT_FONT

    key = font_name.lower().strip()
    return GOOGLE_FONTS_MAP.get(key, _DEFAULT_FONT)


def extract_fonts(pptx_path):
    """Extract the heading and body fonts from a PPTX theme.

    Returns dict with 'title_font', 'content_font' (original names),
    'title_google', 'content_google' (mapped Google Fonts info).
    """
    ns = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    }

    title_font = ''
    content_font = ''

    try:
        prs = Presentation(pptx_path)
        master = prs.slide_masters[0]

        # Find theme element
        theme_el = None
        for rel in master.part.rels.values():
            if 'theme' in rel.reltype:
                theme_part = rel.target_part
                theme_el = etree.fromstring(theme_part.blob)
                break

        if theme_el is not None:
            font_scheme = theme_el.find('.//a:fontScheme', ns)
            if font_scheme is not None:
                major = font_scheme.find('a:majorFont/a:latin', ns)
                minor = font_scheme.find('a:minorFont/a:latin', ns)
                if major is not None:
                    title_font = major.get('typeface', '')
                if minor is not None:
                    content_font = minor.get('typeface', '')
    except (IndexError, AttributeError):
        pass

    title_google = map_to_google_fonts(title_font)
    content_google = map_to_google_fonts(content_font)

    return {
        'title_font': title_font or 'Calibri',
        'content_font': content_font or 'Calibri',
        'title_google': title_google,
        'content_google': content_google,
    }


# ── Style Extraction ─────────────────────────────────────────────────

def extract_style(pptx_path):
    """Analyze the PPTX to infer style metadata for a Layout.

    Returns a style_metadata dict compatible with Layout.style_metadata.
    """
    # Editorial defaults — used as fallback
    style = {
        'border_radius': '16px',
        'border_radius_sm': '12px',
        'card_padding': '32px 28px',
        'top_bar_height': '4px',
        'title_size': '52px',
        'h1_size': '34px',
        'body_size': '16px',
        'shadow': '0 1px 3px rgba(0,0,0,0.04), 0 8px 24px rgba(0,0,0,0.03)',
        'accent_width': '6px',
    }

    try:
        prs = Presentation(pptx_path)

        # Check slide dimensions
        width = prs.slide_width
        height = prs.slide_height
        aspect = width / height if height else 0

        # Scan title placeholder font sizes from the slide master
        master = prs.slide_masters[0]
        for layout in master.slide_layouts:
            for ph in layout.placeholders:
                if ph.placeholder_format.type is not None:
                    # Type 0 = title, type 1 = body
                    pf = ph.placeholder_format
                    if pf.type == 0 or pf.idx == 0:  # Title
                        if ph.text_frame and ph.text_frame.paragraphs:
                            for para in ph.text_frame.paragraphs:
                                if para.font and para.font.size:
                                    pt = para.font.size.pt
                                    style['title_size'] = f'{int(pt * 1.3)}px'
                                    style['h1_size'] = f'{int(pt * 0.85)}px'
                                    break
                            break
            break  # Only check first layout

    except (IndexError, AttributeError):
        pass

    return style


def choose_base_layout(font_info):
    """Determine which existing template directory best matches the extracted fonts.

    Args:
        font_info: dict from extract_fonts()

    Returns:
        str: slug of the base layout ('editorial', 'bold', or 'elegant')
    """
    title_cat = font_info.get('title_google', {}).get('category', 'sans-serif')

    if title_cat == 'display':
        return 'elegant'
    elif title_cat == 'serif':
        return 'editorial'
    elif title_cat == 'sans-serif':
        # Check if it's a "bold" style sans-serif
        title_name = font_info.get('title_font', '').lower()
        bold_fonts = ['impact', 'oswald', 'franklin gothic', 'gotham', 'montserrat']
        if any(bf in title_name for bf in bold_fonts):
            return 'bold'
        return 'bold'  # Default sans-serif → bold
    else:
        return 'editorial'


# ── Image Extraction ─────────────────────────────────────────────────

def extract_images(pptx_path, min_size_bytes=10240):
    """Extract embedded images from a PPTX file.

    Args:
        pptx_path: Path to the .pptx file
        min_size_bytes: Minimum image size to include (filters out icons/bullets)

    Returns:
        list of dicts: [{'data': bytes, 'filename': str, 'content_type': str}]
    """
    prs = Presentation(pptx_path)
    images = []
    seen_hashes = set()

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.shape_type or not hasattr(shape, 'image'):
                continue
            try:
                img = shape.image
                blob = img.blob
                # Skip small images (icons, bullets)
                if len(blob) < min_size_bytes:
                    continue
                # Deduplicate by content hash
                h = hash(blob)
                if h in seen_hashes:
                    continue
                seen_hashes.add(h)

                content_type = img.content_type
                # Derive filename from content type
                ext_map = {
                    'image/png': '.png',
                    'image/jpeg': '.jpg',
                    'image/gif': '.gif',
                    'image/webp': '.webp',
                    'image/tiff': '.tiff',
                    'image/bmp': '.bmp',
                }
                ext = ext_map.get(content_type, '.png')
                filename = f'imported_image_{len(images) + 1}{ext}'

                images.append({
                    'data': blob,
                    'filename': filename,
                    'content_type': content_type,
                })
            except (AttributeError, ValueError):
                continue

    return images


# ── Full Extraction Convenience Function ─────────────────────────────

def extract_all(pptx_path):
    """Run all extractors on a PPTX file.

    Returns dict with 'colors', 'fonts', 'style', 'base_layout', 'images'.
    """
    colors = extract_colors(pptx_path)
    fonts = extract_fonts(pptx_path)
    style = extract_style(pptx_path)
    base_layout = choose_base_layout(fonts)
    images = extract_images(pptx_path)

    return {
        'colors': colors,
        'fonts': fonts,
        'style': style,
        'base_layout': base_layout,
        'images': images,
    }
